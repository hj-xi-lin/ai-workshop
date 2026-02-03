#!/usr/bin/env python3
"""
Customer Performance Report Generator
Generates PowerPoint reports for any customer account with REACH/HIRE split.

Usage:
    python customer_report_generator.py <account_name> [months] [template_path] [output_path]

Arguments:
    account_name: Ultimate parent company name to search for
    months: Number of months to analyze (default: 12)
    template_path: Path to HeyJobs template (default: ~/Desktop/Template.pptx)
    output_path: Output file path (default: ~/Desktop/<account>_Customer_Report.pptx)
"""

import sys
import os
import json
from datetime import datetime, timedelta

# Report data structure (to be populated by Claude via Redshift queries)
REPORT_DATA = {
    "account_name": "",
    "time_range_months": 12,
    "date_from": "",
    "date_to": "",

    # Total metrics
    "total": {
        "booked_budget": 0,
        "used_budget": 0,
        "available_budget": 0,
        "total_jobs": 0,
        "page_views": 0,
        "applications_started": 0,
        "applications_sent": 0,
        "avg_cpa": 0,
    },

    # REACH metrics
    "reach": {
        "booked_budget": 0,
        "used_budget": 0,
        "available_budget": 0,
        "total_jobs": 0,
        "page_views": 0,
        "applications_started": 0,
        "applications_sent": 0,
        "avg_cpa": 0,
        "conversion_rate": 0,
        "monthly_cpa": [],  # [(month_label, cpa, jobs), ...]
        "monthly_apps": [],  # [(month_label, started, sent), ...]
    },

    # HIRE metrics
    "hire": {
        "booked_budget": 0,
        "used_budget": 0,
        "available_budget": 0,
        "total_jobs": 0,
        "page_views": 0,
        "applications_started": 0,
        "applications_sent": 0,
        "avg_cpa": 0,
        "conversion_rate": 0,
        "monthly_cpa": [],
        "monthly_apps": [],
    },
}


def generate_report(data, template_path, output_path):
    """Generate PowerPoint report from data."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    # Load template
    prs = Presentation(template_path)

    # HeyJobs brand colors
    HEYJOBS_PURPLE = RGBColor(102, 45, 145)
    HEYJOBS_DARK = RGBColor(51, 51, 51)
    HEYJOBS_GRAY = RGBColor(128, 128, 128)
    HEYJOBS_CORAL = RGBColor(255, 87, 51)
    HEYJOBS_TEAL = RGBColor(0, 191, 165)

    # Light colors for minimalistic look
    LIGHT_PURPLE = RGBColor(180, 160, 200)
    LIGHT_BLUE = RGBColor(160, 180, 210)
    LIGHT_GREEN = RGBColor(160, 210, 180)
    LIGHT_ORANGE = RGBColor(240, 180, 140)
    LIGHT_TEAL = RGBColor(150, 210, 200)
    LIGHT_CORAL = RGBColor(240, 170, 160)

    # Get slide layouts
    LAYOUT_TITLE = prs.slide_layouts[6]
    LAYOUT_SECTION = prs.slide_layouts[2]
    LAYOUT_CONTENT = prs.slide_layouts[3]

    # Delete all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=HEYJOBS_DARK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return txBox

    def add_kpi_box(slide, left, top, width, height, value, label, bg_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        add_textbox(slide, left, top + 0.15, width, 0.6, value,
                    font_size=24, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
        add_textbox(slide, left, top + 0.65, width, 0.4, label,
                    font_size=10, bold=False, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)

    def add_table(slide, left, top, width, height, table_data, header_color=HEYJOBS_PURPLE):
        rows = len(table_data)
        cols = len(table_data[0])
        table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
        col_width = Inches(width / cols)
        for i in range(cols):
            table.columns[i].width = col_width
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_color
                    for para in cell.text_frame.paragraphs:
                        para.font.bold = True
                        para.font.color.rgb = RGBColor(255, 255, 255)
                        para.font.size = Pt(9)
                        para.alignment = PP_ALIGN.CENTER
                else:
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(8)
                        para.font.color.rgb = HEYJOBS_DARK
                        para.alignment = PP_ALIGN.CENTER
        return table

    def add_line_chart(slide, left, top, width, height, categories, series_data, colors):
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series_data:
            chart_data.add_series(name, values)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
        ).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.has_minor_gridlines = False
        chart.category_axis.has_major_gridlines = False
        chart.category_axis.has_minor_gridlines = False
        for i, color in enumerate(colors):
            if i < len(chart.series):
                chart.series[i].format.line.color.rgb = color
                chart.series[i].format.line.width = Pt(2.5)
        return chart

    def format_number(num):
        """Format number with dots as thousand separators."""
        if num >= 1000000:
            return f"{num/1000000:.2f}M".replace('.', ',')
        elif num >= 1000:
            return f"{num:,.0f}".replace(',', '.')
        return str(int(num))

    def format_currency(num):
        """Format currency."""
        if num >= 1000000:
            return f"€{num/1000000:.2f}M".replace('.', ',')
        elif num >= 1000:
            return f"€{num:,.0f}".replace(',', '.')
        return f"€{num:.2f}".replace('.', ',')

    # ============================================================
    # SLIDE 1: Title Slide
    # ============================================================
    slide = prs.slides.add_slide(LAYOUT_TITLE)
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.clear()

    add_textbox(slide, 0.5, 1.8, 9, 0.8, f"{data['account_name']} Customer Report",
                font_size=36, bold=True, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT)
    add_textbox(slide, 0.5, 2.6, 9, 0.6, "Ihre HeyJobs-Performance Analyse",
                font_size=18, bold=False, color=HEYJOBS_DARK, align=PP_ALIGN.LEFT)
    add_textbox(slide, 0.5, 3.2, 9, 0.5, f"Zeitraum: {data['date_from']} - {data['date_to']} ({data['time_range_months']} Monate)",
                font_size=12, bold=False, color=HEYJOBS_GRAY, align=PP_ALIGN.LEFT)

    # ============================================================
    # SLIDE 2: Executive Summary (Minimalistic)
    # ============================================================
    slide = prs.slides.add_slide(LAYOUT_CONTENT)

    add_textbox(slide, 0.3, 0.2, 9, 0.5, "Performance Analyse im Überblick",
                font_size=20, bold=True, color=RGBColor(255, 255, 255))
    add_textbox(slide, 0.3, 0.55, 9, 0.3, f"Letzte {data['time_range_months']} Monate (REACH + HIRE)",
                font_size=10, bold=False, color=HEYJOBS_GRAY)

    # KPI Row 1
    add_kpi_box(slide, 0.3, 0.95, 2.2, 0.9, format_currency(data['total']['booked_budget']), "Gebuchtes Budget", LIGHT_PURPLE)
    add_kpi_box(slide, 2.6, 0.95, 2.2, 0.9, format_currency(data['total']['used_budget']), "Genutztes Budget", LIGHT_BLUE)
    add_kpi_box(slide, 4.9, 0.95, 2.2, 0.9, format_currency(data['total']['available_budget']), "Verfügbares Budget", LIGHT_GREEN)
    add_kpi_box(slide, 7.2, 0.95, 2.2, 0.9, f"€{data['total']['avg_cpa']:.2f}".replace('.', ','), "⌀ CPA", LIGHT_ORANGE)

    # KPI Row 2
    add_kpi_box(slide, 0.3, 2.0, 2.2, 0.9, format_number(data['total']['total_jobs']), "Jobs", LIGHT_PURPLE)
    add_kpi_box(slide, 2.6, 2.0, 2.2, 0.9, format_number(data['total']['page_views']), "Stellenaufrufe", LIGHT_BLUE)
    add_kpi_box(slide, 4.9, 2.0, 2.2, 0.9, format_number(data['total']['applications_started']), "Bew. gestartet", LIGHT_BLUE)
    add_kpi_box(slide, 7.2, 2.0, 2.2, 0.9, format_number(data['total']['applications_sent']), "Bew. gesendet", LIGHT_GREEN)

    # Key insights
    add_textbox(slide, 0.3, 3.1, 9, 0.3, "REACH vs HIRE:",
                font_size=11, bold=True, color=RGBColor(255, 255, 255))

    reach = data['reach']
    hire = data['hire']
    insights = f"""REACH: {format_number(reach['total_jobs'])} Jobs  |  CPA €{reach['avg_cpa']:.2f}  |  {format_number(reach['applications_sent'])} Bewerbungen  |  {reach['conversion_rate']:.0f}% Conversion
HIRE: {format_number(hire['total_jobs'])} Jobs  |  CPA €{hire['avg_cpa']:.2f}  |  {format_number(hire['applications_sent'])} Bewerbungen  |  {hire['conversion_rate']:.0f}% Conversion""".replace('.', ',')

    add_textbox(slide, 0.3, 3.4, 9, 1.2, insights,
                font_size=9, bold=False, color=HEYJOBS_DARK)

    # ============================================================
    # SLIDE 3: Section - Budget Overview
    # ============================================================
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para.clear()
    add_textbox(slide, 0.5, 2.2, 9, 0.8, "1. Budget-Übersicht",
                font_size=28, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT)

    # ============================================================
    # SLIDE 4: Budget Details by Product
    # ============================================================
    slide = prs.slides.add_slide(LAYOUT_CONTENT)

    add_textbox(slide, 0.3, 0.2, 9, 0.5, "Budget-Übersicht nach Produkt",
                font_size=18, bold=True, color=RGBColor(255, 255, 255))

    # REACH Budget
    add_textbox(slide, 0.3, 0.7, 4.5, 0.3, "REACH", font_size=14, bold=True, color=HEYJOBS_PURPLE)
    reach_fin = [
        ["Metrik", "Betrag (EUR)"],
        ["Gebuchtes Budget", format_currency(reach['booked_budget'])],
        ["Genutztes Budget", format_currency(reach['used_budget'])],
        ["Verfügbares Budget", format_currency(reach['available_budget'])],
    ]
    add_table(slide, 0.3, 1.0, 4.5, 1.6, reach_fin)

    # HIRE Budget
    add_textbox(slide, 5.0, 0.7, 4.5, 0.3, "HIRE", font_size=14, bold=True, color=LIGHT_ORANGE)
    hire_fin = [
        ["Metrik", "Betrag (EUR)"],
        ["Gebuchtes Budget", format_currency(hire['booked_budget'])],
        ["Genutztes Budget", format_currency(hire['used_budget'])],
        ["Verfügbares Budget", format_currency(hire['available_budget'])],
    ]
    add_table(slide, 5.0, 1.0, 4.5, 1.6, hire_fin, header_color=LIGHT_ORANGE)

    # Performance comparison
    add_textbox(slide, 0.3, 2.8, 9, 0.3, "Performance Vergleich", font_size=14, bold=True, color=RGBColor(255, 255, 255))
    perf_compare = [
        ["Metrik", "REACH", "HIRE"],
        ["Jobs", format_number(reach['total_jobs']), format_number(hire['total_jobs'])],
        ["Stellenaufrufe", format_number(reach['page_views']), format_number(hire['page_views'])],
        ["Bewerbungsstarts", format_number(reach['applications_started']), format_number(hire['applications_started'])],
        ["Gesendete Bewerbungen", format_number(reach['applications_sent']), format_number(hire['applications_sent'])],
        ["⌀ CPA", f"€{reach['avg_cpa']:.2f}".replace('.', ','), f"€{hire['avg_cpa']:.2f}".replace('.', ',')],
    ]
    add_table(slide, 0.3, 3.1, 9.4, 2.2, perf_compare)

    # ============================================================
    # SLIDE 5: Section - REACH Performance
    # ============================================================
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para.clear()
    add_textbox(slide, 0.5, 2.2, 9, 0.8, "2. REACH Performance",
                font_size=28, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT)

    # ============================================================
    # SLIDE 6: REACH CPA Chart
    # ============================================================
    if reach['monthly_cpa']:
        slide = prs.slides.add_slide(LAYOUT_CONTENT)

        add_textbox(slide, 0.3, 0.15, 7, 0.4, "REACH: CPA Trending",
                    font_size=16, bold=True, color=RGBColor(255, 255, 255))

        # Average CPA badge
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(0.1), Inches(1.8), Inches(0.45))
        shape.fill.solid()
        shape.fill.fore_color.rgb = HEYJOBS_PURPLE
        shape.line.fill.background()
        add_textbox(slide, 7.8, 0.15, 1.8, 0.35, f"⌀ CPA: €{reach['avg_cpa']:.2f}".replace('.', ','),
                    font_size=11, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)

        categories = [m[0] for m in reach['monthly_cpa']]
        cpa_values = tuple(m[1] for m in reach['monthly_cpa'])
        jobs_values = tuple(m[2]/1000 for m in reach['monthly_cpa'])  # Convert to thousands

        add_line_chart(slide, 0.3, 0.6, 9.4, 4.3, categories,
                       [('CPA (€)', cpa_values), ('Jobs (Tsd.)', jobs_values)],
                       [HEYJOBS_CORAL, HEYJOBS_PURPLE])

    # ============================================================
    # SLIDE 7: REACH Applications Chart
    # ============================================================
    if reach['monthly_apps']:
        slide = prs.slides.add_slide(LAYOUT_CONTENT)

        add_textbox(slide, 0.3, 0.15, 7, 0.4, "REACH: Bewerbungen Trending",
                    font_size=16, bold=True, color=RGBColor(255, 255, 255))

        categories = [m[0] for m in reach['monthly_apps']]
        starts = tuple(m[1]/1000 for m in reach['monthly_apps'])
        sent = tuple(m[2]/1000 for m in reach['monthly_apps'])

        add_line_chart(slide, 0.3, 0.55, 6.5, 4.3, categories,
                       [('Bewerbungsstarts (Tsd.)', starts), ('Gesendete Bew. (Tsd.)', sent)],
                       [HEYJOBS_PURPLE, HEYJOBS_TEAL])

        # Summary stats
        add_textbox(slide, 7.0, 0.6, 2.8, 0.35, "REACH Summen",
                    font_size=12, bold=True, color=RGBColor(255, 255, 255))
        stats_text = f"""Bew. gestartet:
{format_number(reach['applications_started'])}

Bew. gesendet:
{format_number(reach['applications_sent'])}

⌀ Conversion:
{reach['conversion_rate']:.1f}%""".replace('.', ',')
        add_textbox(slide, 7.0, 1.0, 2.8, 4.0, stats_text, font_size=10, bold=False, color=HEYJOBS_DARK)

    # ============================================================
    # SLIDE 8: Section - HIRE Performance
    # ============================================================
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para.clear()
    add_textbox(slide, 0.5, 2.2, 9, 0.8, "3. HIRE Performance",
                font_size=28, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT)

    # ============================================================
    # SLIDE 9: HIRE CPA Chart
    # ============================================================
    if hire['monthly_cpa']:
        slide = prs.slides.add_slide(LAYOUT_CONTENT)

        add_textbox(slide, 0.3, 0.15, 7, 0.4, "HIRE: CPA Trending",
                    font_size=16, bold=True, color=RGBColor(255, 255, 255))

        # Average CPA badge
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(0.1), Inches(1.8), Inches(0.45))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_ORANGE
        shape.line.fill.background()
        add_textbox(slide, 7.8, 0.15, 1.8, 0.35, f"⌀ CPA: €{hire['avg_cpa']:.2f}".replace('.', ','),
                    font_size=11, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)

        categories = [m[0] for m in hire['monthly_cpa']]
        cpa_values = tuple(m[1] for m in hire['monthly_cpa'])
        jobs_values = tuple(m[2] for m in hire['monthly_cpa'])

        add_line_chart(slide, 0.3, 0.6, 9.4, 4.3, categories,
                       [('CPA (€)', cpa_values), ('Jobs', jobs_values)],
                       [HEYJOBS_CORAL, LIGHT_ORANGE])

    # ============================================================
    # SLIDE 10: HIRE Applications Chart
    # ============================================================
    if hire['monthly_apps']:
        slide = prs.slides.add_slide(LAYOUT_CONTENT)

        add_textbox(slide, 0.3, 0.15, 7, 0.4, "HIRE: Bewerbungen Trending",
                    font_size=16, bold=True, color=RGBColor(255, 255, 255))

        categories = [m[0] for m in hire['monthly_apps']]
        starts = tuple(m[1]/1000 for m in hire['monthly_apps'])
        sent = tuple(m[2]/1000 for m in hire['monthly_apps'])

        add_line_chart(slide, 0.3, 0.55, 6.5, 4.3, categories,
                       [('Bewerbungsstarts (Tsd.)', starts), ('Gesendete Bew. (Tsd.)', sent)],
                       [LIGHT_ORANGE, HEYJOBS_TEAL])

        # Summary stats
        add_textbox(slide, 7.0, 0.6, 2.8, 0.35, "HIRE Summen",
                    font_size=12, bold=True, color=RGBColor(255, 255, 255))
        stats_text = f"""Bew. gestartet:
{format_number(hire['applications_started'])}

Bew. gesendet:
{format_number(hire['applications_sent'])}

⌀ Conversion:
{hire['conversion_rate']:.1f}%""".replace('.', ',')
        add_textbox(slide, 7.0, 1.0, 2.8, 4.0, stats_text, font_size=10, bold=False, color=HEYJOBS_DARK)

    # ============================================================
    # SLIDE 11: Section - Performance Summary
    # ============================================================
    slide = prs.slides.add_slide(LAYOUT_SECTION)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para.clear()
    add_textbox(slide, 0.5, 2.2, 9, 0.8, "4. Performance Zusammenfassung",
                font_size=28, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT)

    # ============================================================
    # SLIDE 12: Performance Summary Dashboard (Minimalistic)
    # ============================================================
    slide = prs.slides.add_slide(LAYOUT_CONTENT)

    add_textbox(slide, 0.3, 0.2, 9, 0.5, "Gesamtperformance - REACH vs HIRE",
                font_size=18, bold=True, color=RGBColor(255, 255, 255))

    # REACH KPIs
    add_textbox(slide, 0.3, 0.6, 4.6, 0.3, "REACH", font_size=13, bold=True, color=LIGHT_PURPLE)
    add_kpi_box(slide, 0.3, 0.9, 2.2, 0.9, format_number(reach['total_jobs']), "Jobs", LIGHT_PURPLE)
    add_kpi_box(slide, 2.6, 0.9, 2.2, 0.9, format_number(reach['applications_sent']), "Bew. gesendet", LIGHT_BLUE)
    add_kpi_box(slide, 0.3, 1.95, 2.2, 0.9, f"{reach['conversion_rate']:.1f}%".replace('.', ','), "Conversion", LIGHT_TEAL)
    add_kpi_box(slide, 2.6, 1.95, 2.2, 0.9, f"€{reach['avg_cpa']:.2f}".replace('.', ','), "⌀ CPA", LIGHT_GREEN)

    # HIRE KPIs
    add_textbox(slide, 5.1, 0.6, 4.6, 0.3, "HIRE", font_size=13, bold=True, color=LIGHT_ORANGE)
    add_kpi_box(slide, 5.1, 0.9, 2.2, 0.9, format_number(hire['total_jobs']), "Jobs", LIGHT_ORANGE)
    add_kpi_box(slide, 7.4, 0.9, 2.2, 0.9, format_number(hire['applications_sent']), "Bew. gesendet", LIGHT_CORAL)
    add_kpi_box(slide, 5.1, 1.95, 2.2, 0.9, f"{hire['conversion_rate']:.1f}%".replace('.', ','), "Conversion", LIGHT_CORAL)
    add_kpi_box(slide, 7.4, 1.95, 2.2, 0.9, f"€{hire['avg_cpa']:.2f}".replace('.', ','), "⌀ CPA", LIGHT_ORANGE)

    # Footer
    add_textbox(slide, 0.3, 3.1, 9.4, 0.3,
                f"Zeitraum: {data['date_from']} - {data['date_to']} | Account: {data['account_name']}",
                font_size=8, bold=False, color=HEYJOBS_GRAY, align=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 13: Closing Slide
    # ============================================================
    slide = prs.slides.add_slide(LAYOUT_CONTENT)

    add_textbox(slide, 0.5, 1.5, 9, 0.8, "Gemeinsam die Zukunft Ihres Recruitings gestalten",
                font_size=24, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 2.5, 9, 0.5, "Einfach. Erfolgreich. Einstellen.",
                font_size=16, bold=False, color=HEYJOBS_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 3.5, 9, 0.3, "Vielen Dank!",
                font_size=20, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)

    # Save presentation
    prs.save(output_path)
    print(f"Report saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    # This script is meant to be called with data populated by Claude
    # Example usage for testing:
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    # Load data from JSON file if provided
    data_file = sys.argv[1]
    with open(data_file, 'r') as f:
        data = json.load(f)

    template_path = sys.argv[2] if len(sys.argv) > 2 else os.path.expanduser("~/Desktop/Template.pptx")
    output_path = sys.argv[3] if len(sys.argv) > 3 else os.path.expanduser(f"~/Desktop/{data['account_name']}_Customer_Report.pptx")

    generate_report(data, template_path, output_path)
